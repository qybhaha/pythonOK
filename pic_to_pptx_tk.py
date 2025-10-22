import tkinter as tk
from tkinter import filedialog, messagebox
import os
from pptx import Presentation
from pptx.util import Inches

# 16:9 幻灯片的标准尺寸（单位：英寸）
WIDTH_16_9 = 13.33  # 约为 1920px / 96dpi
HEIGHT_16_9 = 7.5   # 约为 1080px / 96dpi


class ImageListApp:
    def __init__(self, root):
        self.root = root
        self.root.title("图片转PPT工具")

        self.image_files = []
        self.ppt_file_path = ""

        # 创建列表框
        self.image_listbox = tk.Listbox(root, width=60, height=20)
        self.image_listbox.pack(pady=10)

        # 创建按钮
        button_frame = tk.Frame(root)
        button_frame.pack(pady=10)

        load_button = tk.Button(button_frame, text="加载图片", command=self.on_load_images)
        load_button.pack(side=tk.LEFT, padx=5)

        create_ppt_button = tk.Button(button_frame, text="制作PPT", command=self.on_create_ppt)
        create_ppt_button.pack(side=tk.LEFT, padx=5)

        clear_button = tk.Button(button_frame, text="清空图片", command=self.on_clear_images)
        clear_button.pack(side=tk.LEFT, padx=5)

        open_ppt_dir_button = tk.Button(button_frame, text="打开目录", command=self.on_open_ppt_dir)
        open_ppt_dir_button.pack(side=tk.LEFT, padx=5)

    def on_load_images(self):
        filetypes = [("图片文件", "*.jpg;*.png")]
        selected_files = filedialog.askopenfilenames(filetypes=filetypes)
        if selected_files:
            self.image_files = selected_files
            self.image_listbox.delete(0, tk.END)
            for file in self.image_files:
                self.image_listbox.insert(tk.END, file)

    def on_create_ppt(self):
        if not self.image_files:
            messagebox.showinfo("提示", "请先加载图片")
            return

        prs = Presentation()
        # 设置幻灯片尺寸为 16:9
        prs.slide_width = Inches(WIDTH_16_9)
        prs.slide_height = Inches(HEIGHT_16_9)

        for image_file in self.image_files:
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # 添加图片
            pic = slide.shapes.add_picture(image_file, 0, 0)

            # 获取图片和幻灯片的尺寸
            image_width, image_height = pic.image.size
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # 计算图片的宽高比和幻灯片的宽高比
            image_aspect_ratio = image_width / image_height
            slide_aspect_ratio = slide_width / slide_height

            if image_aspect_ratio > slide_aspect_ratio:
                # 图片更宽，裁剪左右部分
                target_height = slide_height
                target_width = target_height * slide_aspect_ratio
                scale = target_height / image_height
                crop_left = (image_width - (target_width / scale)) / 2 / image_width
                crop_right = crop_left
                crop_top = 0
                crop_bottom = 0
            else:
                # 图片更高，裁剪上下部分
                target_width = slide_width
                target_height = target_width / slide_aspect_ratio
                scale = target_width / image_width
                crop_top = (image_height - (target_height / scale)) / 2 / image_height
                crop_bottom = crop_top
                crop_left = 0
                crop_right = 0

            # 设置裁剪
            pic.crop_left = crop_left
            pic.crop_right = crop_right
            pic.crop_top = crop_top
            pic.crop_bottom = crop_bottom

            # 设置图片尺寸为幻灯片尺寸
            pic.width = slide_width
            pic.height = slide_height
            pic.left = 0
            pic.top = 0

        # 保存 PPT
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Files", "*.pptx")])
        if save_path:
            self.ppt_file_path = save_path
            prs.save(self.ppt_file_path)
            messagebox.showinfo("提示", f"PPT 已成功保存到 {self.ppt_file_path}")

    def on_clear_images(self):
        self.image_files = []
        self.image_listbox.delete(0, tk.END)

    def on_open_ppt_dir(self):
        if not self.ppt_file_path:
            messagebox.showinfo("提示", "请先制作PPT")
            return
        ppt_dir = os.path.dirname(self.ppt_file_path)
        if os.name == 'nt':  # Windows
            os.startfile(ppt_dir)
        elif os.name == 'posix':  # Linux or macOS
            os.system(f'open "{ppt_dir}"')


if __name__ == "__main__":
    root = tk.Tk()
    app = ImageListApp(root)
    root.mainloop()
