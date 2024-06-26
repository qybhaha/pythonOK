import wx
from wx.lib.pubsub import pub
from pptx import Presentation
from pptx.util import Inches

class ImageListFrame(wx.Frame):
    def __init__(self, parent, title):
        super(ImageListFrame, self).__init__(parent, title=title, size=(600, 400))

        self.panel = wx.Panel(self)
        self.image_listbox = wx.ListBox(self.panel, size=(500, 500), style=wx.LB_SINGLE)
        self.load_button = wx.Button(self.panel, label='加载图片')
        self.create_ppt_button = wx.Button(self.panel, label='制作PPT')

        self.Bind(wx.EVT_BUTTON, self.on_load_images, self.load_button)
        self.Bind(wx.EVT_BUTTON, self.on_create_ppt, self.create_ppt_button)

        sizer = wx.BoxSizer(wx.HORIZONTAL)
        sizer.Add(self.image_listbox, 1, wx.EXPAND | wx.ALL, 10)
        button_sizer = wx.BoxSizer(wx.VERTICAL)
        button_sizer.Add(self.load_button, 0, wx.ALL, 10)
        button_sizer.Add(self.create_ppt_button, 0, wx.ALL, 10)
        sizer.Add(button_sizer, 0, wx.ALL, 10)
        self.panel.SetSizer(sizer)

        self.image_files = []

    def on_load_images(self, event):
        wildcard = "Image files (*.jpg;*.png)|*.jpg;*.png"
        dialog = wx.FileDialog(self, "Choose images to load", wildcard=wildcard,
                               style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST | wx.FD_MULTIPLE)
        if dialog.ShowModal() == wx.ID_OK:
            self.image_files = dialog.GetPaths()
            self.image_listbox.Clear()
            for file in self.image_files:
                self.image_listbox.Append(file)
        dialog.Destroy()

    def on_create_ppt(self, event):
        if not self.image_files:
            wx.MessageBox("请先加载图片")
            return

        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        blank_slide_layout = prs.slide_layouts[6]

        for idx, image_file in enumerate(self.image_files):
            slide = prs.slides.add_slide(blank_slide_layout)
            left = Inches(0)  # 图片左边距
            top = Inches(0)  # 图片顶边距
            width = Inches(16)  # 图片宽度
            height = Inches(9)  # 图片高度
            slide.shapes.add_picture(image_file, left, top, width=width, height=height)

        save_dialog = wx.FileDialog(self, "Save PowerPoint file", wildcard="PowerPoint files (*.pptx)|*.pptx",
                                   style=wx.FD_SAVE | wx.FD_OVERWRITE_PROMPT)
        if save_dialog.ShowModal() == wx.ID_OK:
            prs.save(save_dialog.GetPath())
            wx.MessageBox(f"PowerPoint file saved as {save_dialog.GetPath()}")
        save_dialog.Destroy()

class ImageListApp(wx.App):
    def OnInit(self):
        self.frame = ImageListFrame(None, title="16:9比例PPTX生成器")
        self.frame.Show()
        return True

if __name__ == '__main__':
    app = ImageListApp()
    app.MainLoop()
